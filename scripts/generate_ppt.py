"""
Generate PRAHARI submission PowerPoint — plain design, maximum substance.
Run: python scripts/generate_ppt.py
Output: PRAHARI_Flipkart_Gridlock_2.0.pptx
"""
import json
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

STATS_PATH = ROOT / "scripts" / "ppt_stats.json"
OUTPUT = ROOT / "PRAHARI_Flipkart_Gridlock_2.0.pptx"

# Plain palette
BLACK = RGBColor(0x1A, 0x1A, 0x1A)
DARK = RGBColor(0x33, 0x33, 0x33)
ACCENT = RGBColor(0x1A, 0x73, 0xE8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x66, 0x66, 0x66)


def load_stats():
    if STATS_PATH.exists():
        return json.loads(STATS_PATH.read_text(encoding="utf-8-sig"))
    return {}


def set_slide_white(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def add_title_bar(slide, title: str, subtitle: str = ""):
    set_slide_white(slide)
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9.0), Inches(0.7))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = BLACK
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.5), Inches(1.05), Inches(9.0), Inches(0.45))
        sp = sub.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(14)
        sp.font.color.rgb = GRAY
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.45), Inches(9.0), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()


def add_bullets(slide, items, top=1.7, left=0.55, width=8.9, font_size=16, line_space=1.15):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(item, tuple):
            text, bold = item
            p.text = text
            p.font.bold = bold
        else:
            p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK
        p.space_after = Pt(6)
        p.line_spacing = line_space
        p.level = 0


def add_section_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_white(slide)
    tbox = slide.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(8.4), Inches(1.2))
    tp = tbox.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(36)
    tp.font.bold = True
    tp.font.color.rgb = ACCENT
    if subtitle:
        sbox = slide.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(8.4), Inches(0.8))
        sp = sbox.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(18)
        sp.font.color.rgb = GRAY


def add_content_slide(prs, title, bullets, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, title, subtitle)
    add_bullets(slide, bullets)
    return slide


def add_two_col_slide(prs, title, left_title, left_items, right_title, right_items, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, title, subtitle)

    lt = slide.shapes.add_textbox(Inches(0.5), Inches(1.65), Inches(4.3), Inches(0.35))
    lt.text_frame.paragraphs[0].text = left_title
    lt.text_frame.paragraphs[0].font.bold = True
    lt.text_frame.paragraphs[0].font.size = Pt(15)
    lt.text_frame.paragraphs[0].font.color.rgb = ACCENT

    rt = slide.shapes.add_textbox(Inches(5.0), Inches(1.65), Inches(4.3), Inches(0.35))
    rt.text_frame.paragraphs[0].text = right_title
    rt.text_frame.paragraphs[0].font.bold = True
    rt.text_frame.paragraphs[0].font.size = Pt(15)
    rt.text_frame.paragraphs[0].font.color.rgb = ACCENT

    add_bullets(slide, left_items, top=2.0, left=0.5, width=4.2, font_size=14)
    add_bullets(slide, right_items, top=2.0, left=5.0, width=4.2, font_size=14)


def build_ppt():
    s = load_stats()
    flow = s.get("flow", {})
    plan = s.get("plan", {})
    rj = s.get("replay_junction", {})
    rh = s.get("replay_hex", {})
    mm = s.get("model_metrics", {})

    off_grid = s.get("off_grid", 147880)
    total = s.get("total_violations", 298450)
    off_pct = off_grid / total * 100 if total else 0

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ── TITLE ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_white(slide)
    t = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(8.4), Inches(1.0))
    t.text_frame.paragraphs[0].text = "PRAHARI"
    t.text_frame.paragraphs[0].font.size = Pt(48)
    t.text_frame.paragraphs[0].font.bold = True
    t.text_frame.paragraphs[0].font.color.rgb = ACCENT

    s1 = slide.shapes.add_textbox(Inches(0.8), Inches(2.9), Inches(8.4), Inches(0.6))
    s1.text_frame.paragraphs[0].text = "Parking Intelligence Command System"
    s1.text_frame.paragraphs[0].font.size = Pt(22)
    s1.text_frame.paragraphs[0].font.color.rgb = DARK

    s2 = slide.shapes.add_textbox(Inches(0.8), Inches(3.7), Inches(8.4), Inches(1.2))
    tf = s2.text_frame
    lines = [
        "Flipkart Gridlock 2.0 — Theme 1: Poor Visibility on Parking-Induced Congestion",
        "Built for Bengaluru Traffic Police (BTP) · ASTraM Violation Data",
        f"Data: {s.get('date_start', 'Nov 2023')} to {s.get('date_end', 'Apr 2024')} · {total:,} parking violations analysed",
    ]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.color.rgb = GRAY

    s3 = slide.shapes.add_textbox(Inches(0.8), Inches(5.8), Inches(8.4), Inches(0.5))
    s3.text_frame.paragraphs[0].text = "[Team Name / Members / Institution — fill before submission]"
    s3.text_frame.paragraphs[0].font.size = Pt(12)
    s3.text_frame.paragraphs[0].font.color.rgb = GRAY

    # ── EXECUTIVE SUMMARY ──
    add_content_slide(prs, "Executive Summary", [
        "PRAHARI is an AI-driven parking enforcement command system that answers one question BTP cannot answer today: where should officers go next — including places with no junction code, no camera, and no patrol assignment.",
        f"The data proves the gap: {off_grid:,} of {total:,} parking violations ({off_pct:.1f}%) occur at locations BTP's junction-based system cannot see.",
        "PRAHARI converts raw violation records into four actionable outputs: (1) blind-spot maps, (2) traffic-flow impact estimates, (3) shift-aware patrol routes, (4) printable field briefing cards.",
        f"Held-out validation (last 14 days, same officer budget): +{rj.get('improvement_pct', 1.4):.1f}% PCIS at junction resolution; +{rh.get('improvement_pct', 5.2):.1f}% when deploying to off-grid H3 hex cells.",
        f"Traffic impact (estimation model): ~{flow.get('daily_avg_delay_veh_hours', 5075):,.0f} vehicle-hours of delay per day from illegal parking alone.",
        "Deliverable: live Streamlit Command Center + full analytical pipeline — deployable on existing BTP data, no new sensors required.",
    ])

    # ── PROBLEM ──
    add_section_slide(prs, "1. The Problem", "Why parking-induced congestion is invisible to current systems")

    add_content_slide(prs, "Problem Statement — Theme 1", [
        "Gridlock 2.0 Theme 1 asks teams to address poor visibility on parking-induced congestion — quantifying impact and improving enforcement where the system is blind.",
        "Bengaluru's ASTraM platform captures ~298K parking violations over 6 months — but visibility is structurally uneven.",
        "Traffic AI cameras excel at junction-coded violations (87%+ of non-parking traffic). Parking enforcement remains manual: ~1,500 challans/day by officers on foot.",
        "Officers patrol by experience and known junction lists. Locations without junction codes are invisible to routing, reporting, and camera systems.",
        "Result: chronic illegal parking persists in residential streets, commercial corridors, and side roads — blocking carriageway, delaying traffic, and eroding public trust.",
    ], "Poor Visibility on Parking-Induced Congestion")

    add_two_col_slide(prs, "What BTP Sees vs. What BTP Misses",
        "Known System (Junction-Coded)",
        [
            f"{s.get('with_junction', 150570):,} violations at {s.get('junctions', 169)} named junctions",
            "Mapped in ASTraM with junction codes (BTP051, BTP063, etc.)",
            "Eligible for camera enforcement and junction-level reporting",
            "Officers can be assigned to named hotspots",
            "Still under-patrolled relative to violation density in many zones",
        ],
        "Invisible System (Off-Grid)",
        [
            f"{off_grid:,} violations with junction_name = 'No Junction'",
            f"{off_pct:.1f}% of all parking violations — nearly half the dataset",
            "No junction code → no camera trigger, no patrol assignment",
            "HDBSCAN discovers 1,400+ hidden hotspot clusters in this space",
            "This is where congestion compounds silently",
        ],
    )

    add_content_slide(prs, "Why This Matters for Bengaluru", [
        "Parking violations are not random noise — they physically block carriageway and reduce road capacity.",
        f"Estimated impact (IRC/HCM-based model): {flow.get('total_carriageway_blocked_m2', 2945080):,.0f} m² carriageway blocked over 6 months.",
        f"Daily average: ~{flow.get('daily_avg_delay_veh_hours', 5075):,.0f} vehicle-hours of delay; ~{flow.get('daily_avg_delay_person_hours', 7612):,.0f} person-hours.",
        f"~{flow.get('junction_delay_pct', 49):.0f}% of delay originates at junction-coded locations; the rest is off-grid.",
        f"Peak enforcement hour (IST): {s.get('peak_hour', 10)}:00 — morning commercial activity window.",
        f"Estimated uncollected fines (₹500/violation): ₹{flow.get('est_total_uncollected_fines_rs', 149225000)/10000000:.1f} Cr over 6 months — illustrative, not audited revenue.",
        "Without spatial intelligence beyond junction codes, BTP cannot prioritise the highest-impact locations each shift.",
    ])

    # ── INSIGHT ──
    add_section_slide(prs, "2. Our Insight", "The enforcement visibility gap is the root cause")

    add_content_slide(prs, "Core Thesis", [
        "BTP does not have a data shortage. BTP has a visibility and deployment gap.",
        "298K records already contain GPS coordinates, vehicle type, violation type, officer ID, and timestamp — enough to build a full intelligence layer.",
        "The missing layer is not more cameras — it is a system that: (a) discovers off-grid hotspots, (b) scores impact consistently, (c) deploys officers shift-by-shift, (d) proves improvement on held-out data.",
        "PRAHARI treats parking enforcement as a command-and-control problem, not a dashboard problem.",
        "Judges should remember one number: 147,880 — violations BTP's junction system is structurally blind to.",
    ])

    add_content_slide(prs, "Design Principles", [
        "Evidence-first: every headline metric is traceable to ASTraM records or a stated estimation model.",
        "Fair comparison: PRAHARI is validated against BTP's actual patrol pattern — same officers, same stop budget, same intercept radius.",
        "No data leakage: prediction model and replay use temporal train/test splits; shift routing uses only causal (prior) data.",
        "Operational, not academic: output is patrol routes and printable field cards — what an inspector runs each morning.",
        "Transparent assumptions: traffic-flow numbers are labeled estimates with IRC/HCM sources — challengeable, not hidden.",
        "Deployable today: runs on existing violation CSV; no new hardware, no Map API dependency for core product.",
    ])

    # ── SOLUTION ──
    add_section_slide(prs, "3. PRAHARI — The Solution", "Parking Intelligence Command System")

    add_content_slide(prs, "Solution Overview", [
        "PRAHARI = end-to-end pipeline from raw ASTraM violations → shift-ready patrol deployment.",
        "Layer 1 — Data Intelligence: clean, IST-normalise, H3-index, score every violation with PCIS.",
        "Layer 2 — Spatial Discovery: HDBSCAN blind-spot clustering + enforcement-exposure bias correction.",
        "Layer 3 — Prediction: LightGBM hurdle model forecasts where violations will occur next.",
        "Layer 4 — Impact Quantification: carriageway blocked, capacity reduction, vehicle-hours delay.",
        "Layer 5 — Optimisation: load-balanced patrol routing under shift time constraints.",
        "Layer 6 — Command Interface: Streamlit Patrol Command Center with field briefing export.",
    ])

    add_content_slide(prs, "System Architecture", [
        "INPUT: BTP ASTraM anonymised violation CSV (Jan–May 2024, 298K parking records)",
        "  ↓  Data Pipeline (pandas + H3): clean · PCIS score · cell aggregation · recurrence",
        "  ↓  Spatial Analysis (HDBSCAN + exposure scoring): blind spots · bias-corrected priority map",
        "  ↓  Prediction (LightGBM hurdle): binary occurrence + intensity regression per H3 cell",
        "  ↓  Traffic Flow Module (IRC/HCM): physical delay and capacity metrics",
        "  ↓  Optimizer (OR-Tools + greedy fallback): multi-unit route assignment",
        "  ↓  OUTPUT: Patrol Command Center · Proof tab · Intelligence tabs · Field cards",
        "Deployment: Streamlit Community Cloud · pre-bundled parquet (~24 MB) · no raw 105 MB CSV needed",
    ])

    add_content_slide(prs, "What Makes PRAHARI Different", [
        "Not another heatmap — a shift deployment engine with stop sequences and arrival windows.",
        "Not junction-only — H3 hex resolution captures off-grid hotspots BTP has no codes for.",
        "Not vanity ML — hurdle model with Precision@10% = 100% on top-ranked cells; features audited for leakage.",
        "Not unfair baselines — replay compares against BTP's actual historical patrol, not random points.",
        "Not hidden physics — traffic delay model shows every assumption on screen for reviewer challenge.",
        "Not prototype-only — pre-computed bundles enable instant landing page on 1 GB cloud instances.",
    ])

    # ── DATA ──
    add_section_slide(prs, "4. Data & Methodology", "298,450 parking violations · 151 days · Bengaluru")

    add_content_slide(prs, "Dataset Overview", [
        f"Source: BTP ASTraM anonymised police violation dataset (Hackathon release)",
        f"Period: {s.get('date_start')} to {s.get('date_end')} ({flow.get('data_days', 151)} unique days)",
        f"Parking violations extracted: {total:,} (filtered from mixed violation types via keyword matching)",
        f"Geography: {s.get('stations', 54)} police station jurisdictions · {s.get('junctions', 169)} unique junction names",
        "Fields used: GPS (lat/lng), violation_type, vehicle_type, junction_name, police_station, created_datetime, device_id, officer_id",
        "Temporal normalisation: UTC → Asia/Kolkata (IST) for all hour/day/shift analysis",
        "Spatial indexing: Uber H3 resolution 9 (~174 m edge, ~0.1 km² per hex) for fine analysis; resolution 7 for overview",
    ])

    add_content_slide(prs, "PCIS — Parking Congestion Impact Score", [
        "Problem: raw violation counts treat a parked lorry on a main road the same as a scooter on a side street.",
        "PCIS = vehicle_weight × violation_severity × location_weight — a unified severity score per event.",
        "Vehicle weight (examples): HGV/Lorry = 5.0 · Bus = 4.0 · Car = 2.0 · Two-wheeler = 1.0",
        "Violation severity (examples): Main road parking / double parking = 3.0 · No parking = 2.0 · Footpath = 1.5",
        "Location weight: named junction = 2.0 · main-road violation off junction = 1.5 · other = 1.0",
        "PCIS enables fair ranking across vehicle types, violation types, and locations — used in routing, replay, and reporting.",
        "All weights defined in config.py — auditable and adjustable by BTP policy.",
    ])

    add_content_slide(prs, "Data Pipeline", [
        "Step 1: Parse violation_type JSON arrays; filter to parking-related keywords (14 violation types)",
        "Step 2: Resolve vehicle type (updated_vehicle_type fallback to vehicle_type)",
        "Step 3: Convert timestamps to IST; derive hour, day-of-week, 4-hour time windows, weekend flag",
        "Step 4: Assign H3 indices at resolutions 7 and 9 for every GPS point",
        "Step 5: Compute PCIS per record; aggregate to (H3 cell × date × time_window) grids",
        "Step 6: Compute recurrence factors per cell/time-window for prediction features",
        "Step 7: Build junction summary (deduplicated by junction_name — fixes 135 duplicate rows from multi-station mapping)",
        "Output: parking.parquet + cell_agg + recurrence + junction_summary (~24 MB bundled for cloud deploy)",
    ])

    # ── SPATIAL ──
    add_section_slide(prs, "5. Spatial Intelligence", "Finding what junction codes cannot see")

    add_content_slide(prs, "Blind Spot Discovery — HDBSCAN", [
        f"Input: {off_grid:,} violations where junction_name = 'No Junction'",
        "Algorithm: HDBSCAN density clustering on haversine distance (lat/lng in radians)",
        "Parameters: min_cluster_size = 20 · min_samples = 10",
        f"Output: 1,436 discovered hidden hotspot clusters ranked by total PCIS",
        "Each cluster profiled: violation count, days active, dominant vehicle, heavy-vehicle ratio, main-road ratio, peak hour, police station jurisdiction",
        "Example autopsy (worst cluster): thousands of violations over 6 months, active most days, zero junction code, zero camera, zero patrol assignment",
        "PRAHARI names and maps these clusters — making the invisible actionable for the first time.",
    ])

    add_content_slide(prs, "Enforcement Exposure & Bias Correction", [
        "Raw violation density alone is misleading — some cells are over-patrolled, others never visited.",
        "Exposure metric: unique_officers × unique_dates per H3 cell (from ASTraM officer/device IDs)",
        "Blind spot flag: low exposure (bottom quartile) AND high PCIS-per-exposure (top quartile)",
        "Over-patrolled flag: high exposure AND low total PCIS — resources wasted",
        "Bias-corrected score: PCIS_per_exposure × log(1 + violation_count) — prioritises under-enforced high-impact cells",
        "Displacement analysis module tracks whether enforcement at one cell shifts violations to neighbours",
        "This prevents PRAHARI from simply recommending places BTP already saturates with patrols.",
    ])

    add_two_col_slide(prs, "Spatial Resolution Comparison",
        "Junction-Level (BTP Today)",
        [
            "169 named junctions",
            "Works for coded hotspots",
            "Misses all off-grid violations",
            "Replay improvement: +1.4% PCIS",
            "Fair test — same budget both sides",
        ],
        "H3 Hex-Level (PRAHARI)",
        [
            "Thousands of H3 resolution-9 cells",
            "Deploys to off-grid hotspots",
            "Captures 147,880 invisible violations",
            "Replay improvement: +5.2% PCIS",
            "99.8% of hex gain from off-grid interceptions",
        ],
        "Why hex resolution matters for Theme 1",
    )

    # ── PREDICTION ──
    add_section_slide(prs, "6. Predictive Intelligence", "Where will violations occur next?")

    add_content_slide(prs, "LightGBM Hurdle Model", [
        "Two-stage model (standard for sparse count data):",
        "  Stage 1 — Classifier: will this H3 cell have any violation in this time window? (binary)",
        "  Stage 2 — Regressor: if yes, predict PCIS intensity (continuous)",
        "Training cutoff: 2024-02-15 · Test: after cutoff (temporal split — no future leakage)",
        "Features (14, all causal): time_window, dow, weekend, month, hour sin/cos, dow sin/cos, lag-1/6/42, rolling means, recurrence_factor",
        "Explicitly excluded (leakage audit): heavy_ratio, main_road_ratio — these are outcomes, not predictors",
        "Top 500 active H3 cells · grid expanded to include zero-violation slots for proper hurdle training",
    ])

    add_content_slide(prs, "Model Performance", [
        f"Precision@5%:  {mm.get('precision_at_5pct', 0.61)*100:.1f}%  ·  Hit Rate@5%:  {mm.get('hit_rate_at_5pct', 0.96)*100:.1f}%",
        f"Precision@10%: {mm.get('precision_at_10pct', 1.0)*100:.1f}% ·  Hit Rate@10%: {mm.get('hit_rate_at_10pct', 0.88)*100:.1f}%",
        f"Precision@20%: {mm.get('precision_at_20pct', 1.0)*100:.1f}% ·  Hit Rate@20%: {mm.get('hit_rate_at_20pct', 0.90)*100:.1f}%",
        f"Binary classifier F1: {mm.get('f1', 0.61)*100:.1f}% · Precision: {mm.get('precision', 0.67)*100:.1f}% · Recall: {mm.get('recall', 0.57)*100:.1f}%",
        f"Intensity MAE: {mm.get('mae', 2.38):.2f} PCIS · RMSE: {mm.get('rmse', 12.16):.1f}",
        "Interpretation: when PRAHARI flags the top 10% of predicted cells, 100% actually had violations — officers won't be sent to empty locations.",
        "Top features: rolling means, lag features, recurrence factor, temporal encodings — violations are habitual and time-patterned.",
    ])

    add_content_slide(prs, "Enforcement Efficacy Analysis", [
        "Observational before/after: do violations decrease at junctions following high-enforcement weeks?",
        "Method: identify enforcement spikes (top quartile officer activity per junction-week); compare PCIS in following 2 weeks vs prior 2 weeks",
        "Caveat stated on every screen: correlation, not causal proof — regression to mean possible",
        "Identifies responsive junctions (enforcement worked) vs resistant junctions (violations persisted/increased)",
        "Feeds Command Center rationale: some locations need different tactics (towing, signage) not just more patrols",
        "Chronic offender analysis and station-level comparison modules included in System Health tab",
    ])

    # ── TRAFFIC FLOW ──
    add_section_slide(prs, "7. Traffic Flow Impact", "Quantifying congestion — as Theme 1 requires")

    add_content_slide(prs, "From Violations to Physical Metrics", [
        "Theme 1 explicitly asks to quantify parking impact on traffic flow. PCIS alone is relative — PRAHARI converts to physical units.",
        "Per-violation calculations:",
        "  • Carriageway blocked (m²) = (vehicle length + buffer) × (vehicle width + buffer)",
        "  • Capacity reduction (%) = effective width reduction / road width — geometric lower bound",
        "  • Vehicle-hours delay = affected flow rate × parking duration × per-vehicle delay (BPR-style)",
        "Road profiles (IRC:86 urban standards): main road 10.5 m (3 lanes) · junction 10.5 m · residential 6.0 m (2 lanes)",
        "Flow rates: 1,200 PCU/lane/hr (arterial) · 800 PCU/lane/hr (residential) per IRC:106-1990",
        "Assumptions: 2 hr avg parking duration · 8 sec avg delay per passing vehicle · 1.5 persons/vehicle",
    ])

    add_content_slide(prs, "City-Wide Traffic Impact (6 Months)", [
        f"Total carriageway blocked: ~{flow.get('total_carriageway_blocked_m2', 2945080):,.0f} m²",
        f"Average capacity reduction (lower bound): ~{flow.get('avg_capacity_reduction_pct', 26):.0f}% per violation event",
        f"Total vehicle-hours of delay: ~{flow.get('total_delay_vehicle_hours', 766298):,.0f}",
        f"Total person-hours of delay: ~{flow.get('total_delay_person_hours', 1149447):,.0f}",
        f"Daily average vehicle-hours lost: ~{flow.get('daily_avg_delay_veh_hours', 5075):,.0f}",
        f"Daily average person-hours lost: ~{flow.get('daily_avg_delay_person_hours', 7612):,.0f}",
        f"Heavy vehicles account for ~{flow.get('heavy_vehicle_blocked_pct', 18):.0f}% of total road area blocked",
        f"Junction-coded violations contribute ~{flow.get('junction_delay_pct', 49):.0f}% of total delay",
        "Example: one parked lorry on a 3-lane main road → ~43% capacity reduction (lower bound) → ~6.9 vehicle-hours from a single 2-hour event",
        "All figures labeled as estimation model outputs — transparent, challengeable, conservative (lower bounds).",
    ])

    # ── OPTIMIZER ──
    add_section_slide(prs, "8. Patrol Optimisation", "Turning intelligence into deployment")

    add_content_slide(prs, "Patrol Routing Engine", [
        "Problem: given N officers, a 4-hour shift, and ~30 min dwell per stop — which locations maximise PCIS intercepted?",
        "Candidate pool: top 25 hotspots by shift-window PCIS ranking (configurable)",
        "Primary solver: Google OR-Tools VRP with time windows and capacity constraints",
        "Production fallback: load-balanced greedy router — deduplicates junctions, travel-aware, balances workload across units",
        f"Load balance ratio (default 5-unit plan): {plan.get('load_balance', 0.74)*100:.0f}% (min route PCIS / max route PCIS)",
        "Output per unit: ordered stop sequence · PCIS target · route map · printable field card with arrival windows",
        "Scenario simulator: marginal PCIS gain curve vs number of units — identifies optimal deployment (knee point)",
    ])

    add_content_slide(prs, "Shift-Aware Deployment", [
        "Parking patterns change by time of day — a morning plan is wrong for evening.",
        "Four shift windows: Morning (08–12) · Midday (12–16) · Evening (16–20) · Night (20–24)",
        "Each shift re-ranks junctions by PCIS within that hour window only — causal, no future data",
        "Verified: top-3 junctions differ between morning and evening shifts — static daily plans miss this",
        "PRAHARI's replay uses per-shift causal ranking: for each held-out day/shift, only prior data informs deployment",
        "This mirrors how an inspector actually thinks: 'where is it bad RIGHT NOW in this shift?'",
    ])

    add_content_slide(prs, "Baseline Comparison — Honest Numbers", [
        "PRAHARI compared against three baselines on the SAME candidate pool and stop budget:",
        "  1. Naive historical: top-N junctions by violation count, same stop budget — improvement: 0.0% (tie)",
        "  2. Random deployment: random sample from candidate pool (100 simulations) — PRAHARI beats random significantly",
        "  3. Optimised routing: load-balanced multi-unit assignment — wins on PCIS per officer-hour via travel-aware sequencing",
        "Why naive tie is honest and important: BTP already knows top junction hotspots. PRAHARI's junction-level value is shift-timing + load balance + off-grid hex deployment.",
        "The +5.2% replay gain comes almost entirely from deploying to off-grid H3 cells — the visibility gap, not re-sorting known junctions.",
    ])

    # ── PROOF ──
    add_section_slide(prs, "9. Proof It Works", "Held-out enforcement replay — no leakage")

    add_content_slide(prs, "Enforcement Replay Methodology", [
        "Question: if BTP had deployed PRAHARI's plan during the last 2 weeks, would it have intercepted more violations?",
        "Train period: 2023-11-10 to 2024-03-25 · Held-out: 2024-03-26 to 2024-04-08 (14 days, never seen during training)",
        "Both strategies: 5 units × 8 stops = 40 locations per shift · 0.35 km intercept radius · same for both sides",
        "ACTUAL (BTP today): top coded junctions by training-period violation count — fixed list, no shift awareness",
        "PRAHARI: causal ranking per shift window — junction or H3 hex resolution",
        "Scoring: sum PCIS of held-out violations within 0.35 km of a deployed location during the matching shift",
        "Observational replay — measures coverage of where violations occurred, not causal enforcement effect (stated explicitly)",
    ])

    add_two_col_slide(prs, "Replay Results — Held-Out Period",
        f"Junction Resolution (+{rj.get('improvement_pct', 1.4):.1f}%)",
        [
            f"ACTUAL intercepted PCIS: {rj.get('actual_pcis', 38806):,.0f}",
            f"PRAHARI intercepted PCIS: {rj.get('prahari_pcis', 39334):,.0f}",
            "Both limited to coded junctions",
            "Gain from shift-aware re-ranking",
            "Conservative, apples-to-apples test",
        ],
        f"Hex Resolution (+{rh.get('improvement_pct', 5.2):.1f}%)",
        [
            f"ACTUAL intercepted PCIS: {rh.get('actual_pcis', 38806):,.0f}",
            f"PRAHARI intercepted PCIS: {rh.get('prahari_pcis', 40830):,.0f}",
            "PRAHARI deploys to off-grid H3 cells",
            "ACTUAL still limited to junctions",
            "99.8% of gain from off-grid hex interceptions",
        ],
        f"Holdout: {rh.get('holdout', '2024-03-26 to 2024-04-08')}",
    )

    add_content_slide(prs, "Why We Show Both Resolutions", [
        "Junction replay (+1.4%): fair head-to-head — proves shift-aware routing adds value even within BTP's existing junction framework.",
        "Hex replay (+5.2%): shows the full PRAHARI vision — deploying where BTP has no codes. ACTUAL cannot compete here by design.",
        "We do not lead with +5.2% — we lead with 147,880 blind violations, then show hex replay as evidence of what fixing visibility unlocks.",
        "Every number in replay is auditable: exact train/holdout split, stop budget, intercept radius, per-day/per-shift deployment logged in debug panel.",
        "This is the difference between a demo and a defensible evaluation.",
    ])

    # ── PRODUCT ──
    add_section_slide(prs, "10. Patrol Command Center", "What BTP runs each shift")

    add_content_slide(prs, "Command Center — Landing Page", [
        "Designed as a command interface, not an analytics dashboard.",
        "Hero metric: 147,880 off-grid violations — the visibility gap in one number.",
        "Contrast band: 'Today WITHOUT PRAHARI' (blind patrol by experience) vs 'Today WITH PRAHARI' (surgical deployment).",
        "Live patrol plan: 5 units · ordered stop sequences · PCIS targets · route map · load balance stats.",
        "Pre-computed default plan: instant load on cloud (<3 sec) — no 20-second spinner on first impression.",
        "Slider reconfiguration: change units/shift/hotspots → re-optimises with cached results.",
        "One-click export: printable field briefing cards (TXT + HTML) with arrival windows per stop.",
    ])

    add_content_slide(prs, "Dashboard — 8 Sections", [
        "1. Patrol Command Center — shift deployment engine (landing page)",
        "2. Proof It Works — enforcement replay with junction/hex toggle + audit debug panel",
        "3. Intelligence: Enforcement Gap — visible vs invisible violation maps",
        "4. Intelligence: Traffic Flow Impact — carriageway, delay, capacity metrics + full assumptions",
        "5. Intelligence: Hotspot Analysis — junction rankings, hourly patterns, vehicle mix",
        "6. Intelligence: Blind Spots — HDBSCAN cluster autopsy + top-30 hidden hotspots",
        "7. Intelligence: Model & Efficacy — prediction metrics + enforcement responsiveness",
        "8. Intelligence: System Health — rejection rates, chronic offenders, station comparison",
        "Lazy page loading: heavy computation runs only when a section is opened — cloud-safe architecture.",
    ])

    add_content_slide(prs, "Field Deployment Workflow", [
        "Morning (Inspector): open Command Center → select shift window → review 5 unit assignments → print field cards.",
        "Each field card: officer name blank · shift · PCIS target · ordered stop sequence · illustrative arrival windows.",
        "Officers patrol assigned sequence with ~30 min dwell per stop; prioritise heavy-vehicle and main-road violations.",
        "End of shift: violations logged in ASTraM as usual — feeds back into next day's intelligence.",
        "Weekly: review Proof tab replay metrics · adjust unit count via scenario simulator · investigate new blind-spot clusters.",
        "No new apps for officers — printable cards work on paper. Streamlit is the inspector's tool.",
    ])

    # ── IMPACT ──
    add_section_slide(prs, "11. Impact & Value for BTP", "Why this wins for Bengaluru")

    add_content_slide(prs, "Operational Impact", [
        f"Visibility: names and maps {off_grid:,} previously invisible violations across 1,400+ hidden clusters",
        "Deployment: shift-specific patrol routes in seconds — replaces experience-based guesswork",
        f"Coverage: +{rh.get('improvement_pct', 5.2):.1f}% more high-impact violations intercepted (held-out hex replay, same budget)",
        f"Balance: load-balanced routes — no officer gets 3× the workload of another (was 69K–226K PCIS spread, now ~108K–124K)",
        "Accountability: every route decision traceable to PCIS ranking + stated assumptions",
        "Field-ready: printable briefing cards — zero training overhead for patrol officers",
    ])

    add_content_slide(prs, "Strategic Impact", [
        "Theme 1 compliance: quantifies parking-induced congestion in vehicle-hours, person-hours, and carriageway m² — not just violation counts.",
        "Resource optimisation: scenario simulator shows diminishing returns — prevents over-deploying officers to saturated zones.",
        "Policy input: identifies resistant junctions where patrol alone doesn't work — signals need for infrastructure changes.",
        "Revenue illustration: ₹149 Cr estimated uncollected fines over 6 months (₹500/violation) — enforcement value, not forecast.",
        "Scalable: pipeline runs on monthly CSV refresh — no custom infrastructure beyond a Streamlit or internal server deploy.",
        "Replicable: H3 + HDBSCAN + hurdle model framework applies to any Indian city with GPS-tagged violation data.",
    ])

    add_content_slide(prs, "Before vs After PRAHARI", [
        ("BEFORE", True),
        "Officers patrol known junctions by experience",
        "147,880 off-grid violations unmapped and unassigned",
        "No shift-specific re-ranking — same plan all day",
        "No traffic-flow quantification for parking",
        "No held-out proof that deployment strategy works",
        "Enforcement bias uncorrected — over/under-patrolled cells invisible",
        "",
        ("AFTER PRAHARI", True),
        "Shift-aware routes with PCIS-ranked stop sequences",
        "1,436 hidden hotspot clusters discovered and prioritised",
        "Morning ≠ Evening deployment — causal per-shift ranking",
        "~5,075 daily vehicle-hours of delay quantified (estimation model)",
        "+5.2% held-out PCIS improvement at hex resolution (same budget)",
        "Bias-corrected exposure map directs officers to under-enforced cells",
    ])

    # ── TECH ──
    add_section_slide(prs, "12. Technical Implementation", "Built to deploy, not just demo")

    add_content_slide(prs, "Technology Stack", [
        "Language: Python 3.10+",
        "Data: pandas · numpy · pyarrow (parquet bundles)",
        "Spatial: Uber H3 (v4) · HDBSCAN · haversine distance",
        "ML: LightGBM hurdle model · scikit-learn metrics · time-series split",
        "Optimisation: Google OR-Tools VRP + custom load-balanced greedy fallback",
        "Frontend: Streamlit · Plotly · Folium maps",
        "Deployment: Streamlit Community Cloud · pre-bundled data (~24 MB) · requirements.txt pinned",
        "All modules in src/ — modular, testable, no monolith",
    ])

    add_content_slide(prs, "Cloud-Ready Architecture", [
        "Problem: raw CSV is 105 MB — exceeds GitHub/HackerEarth 50 MB limit.",
        "Solution: one-time bundle_data.py → parquet + JSON precomputes (~24 MB total).",
        "Pre-computed at build time: flow impact · blind spots (top 30) · default patrol plan · junction/cell aggregations.",
        "Runtime: load parquet (~2 sec) · instant Command Center landing · lazy page loading for heavy tabs.",
        "Clean venv verified: h3, hdbscan, lightgbm, ortools all install from wheels — no build failures.",
        "Memory-safe: HDBSCAN and model training run only when user opens those tabs — survives 1 GB Streamlit Cloud tier.",
    ])

    add_content_slide(prs, "Data Integrity & Audit Trail", [
        "UTC→IST conversion verified for all temporal analysis (peak hour = 10:00 AM IST)",
        "Junction summary deduplicated: was 270 rows (135 duplicates from multi-station mapping) → 135 unique junctions",
        "Prediction leakage audit: removed heavy_ratio and main_road_ratio from features (post-hoc composition metrics)",
        "Replay fairness audit: same budget · same radius · same candidate pool · causal shift ranking · logged in debug panel",
        "Revenue calc corrected: was stops × 15 × ₹500 → now stops × ₹500 (one challan per stop)",
        "Capacity reduction labeled 'lower bound' — real throughput loss worse due to merging friction",
        "Full audit.py module for reproducible verification of all headline numbers",
    ])

    # ── LIMITATIONS ──
    add_section_slide(prs, "13. Honest Limitations", "What we claim — and what we don't")

    add_content_slide(prs, "Limitations & Caveats", [
        "Traffic flow numbers are estimation model outputs (IRC/HCM assumptions) — not field-measured delay.",
        "Enforcement replay is observational coverage analysis — not a randomised controlled trial of enforcement effect.",
        "Efficacy analysis (before/after enforcement weeks) shows correlation — regression to mean is possible.",
        "Prediction model trained on top 500 cells — extrapolation to new areas requires more data.",
        "HDBSCAN clusters are spatial — they don't account for road network topology (no routing graph).",
        "We do not claim causal fine revenue or exact delay seconds — we claim prioritisation improvement with stated assumptions.",
        "Judges can challenge any assumption in the Traffic Flow tab — that transparency is intentional.",
    ])

    # ── ROADMAP ──
    add_section_slide(prs, "14. Roadmap", "Beyond the hackathon")

    add_content_slide(prs, "Phase 2 — Production Integration", [
        "MapMyIndia API: live road width validation · geocoding for cluster naming · turn-by-turn route distances (architecture designed, not live in demo)",
        "ASTraM real-time feed: replace monthly CSV with daily/hourly pipeline refresh",
        "FastAPI backend: REST endpoints for patrol plan, blind spots, replay — enable mobile app integration",
        "Officer mobile app: receive field cards · mark stops complete · GPS verify presence",
        "Feedback loop: post-enforcement violation drop tracked per cell → auto-calibrate PCIS weights",
        "Multi-city: parameterise road profiles and vehicle weights per city · deploy to Mumbai, Delhi, Chennai",
    ])

    add_content_slide(prs, "Phase 3 — Policy & Infrastructure", [
        "Integrate with BTP towing dispatch for heavy-vehicle clusters",
        "Signal timing impact module: parking near signals → spillback delay (extension of flow model)",
        "Public dashboard: citizen-facing parking congestion heatmap (anonymised)",
        "Policy simulator: 'what if we add no-parking zones here?' — counterfactual PCIS reduction",
        "Integration with BMTC bus lane data for bus-stop parking enforcement prioritisation",
    ])

    # ── DEMO GUIDE ──
    add_section_slide(prs, "15. Demo Script", "3-minute video pitch guide")

    add_content_slide(prs, "Recommended Demo Flow (3 Minutes)", [
        "0:00–0:30 — Command Center landing: '147,880 violations BTP cannot see. Here is today's patrol plan.' Show contrast band + route map.",
        "0:30–1:00 — Blind Spots tab: cluster autopsy — one hidden hotspot, thousands of violations, no junction code. 'This is the visibility gap.'",
        "1:00–1:30 — Traffic Flow tab: '~5,075 vehicle-hours lost daily — here are the assumptions.' Show one lorry example.",
        "1:30–2:00 — Proof tab: toggle Junction (+1.4%) then Hex (+5.2%). 'Same officers, same budget — more impact because we deploy off-grid.'",
        "2:00–2:30 — Change shift slider → plan re-ranks. Download field cards. 'This is what the inspector prints each morning.'",
        "2:30–3:00 — Close: 'No new sensors. Existing data. Shift-ready deployment. PRAHARI — because Bengaluru cannot enforce what it cannot see.'",
    ])

    add_content_slide(prs, "Key Talking Points for Judges", [
        "Lead with the problem (147,880 blind), not the model accuracy.",
        "Show Proof tab before claiming improvement — methodology is on screen.",
        "When asked 'is this causal?': 'No — it's held-out coverage replay with fair baselines. Here is the audit panel.'",
        "When asked 'are delay numbers real?': 'They are IRC-based estimates with every assumption listed. Conservative lower bounds.'",
        "When asked 'why not just add cameras?': 'Cameras need junction codes. 49.5% of violations have none. PRAHARI finds them first.'",
        "When asked 'what does BTP do tomorrow?': 'Print field cards from Command Center. Same officers, better targets, this shift.'",
    ])

    # ── APPENDIX ──
    add_section_slide(prs, "Appendix", "Reference material for evaluators")

    add_content_slide(prs, "A1 — PCIS Weight Tables", [
        "Vehicle Weights: HGV/Lorry/Tanker = 5.0 · Bus variants = 4.0 · Tempo/LGV/Mini Lorry/Tractor = 3.0–3.5",
        "  Car/Jeep/Van/School Vehicle = 2.0 · Auto = 1.5 · Two-wheeler = 1.0 · Default = 1.5",
        "Violation Severity: Main road parking / Double parking / One-way = 3.0 · Near crossing/signal/bus stop = 2.5",
        "  HTV prohibited / Wrong parking / No parking = 2.0 · Footpath / Other bus stop = 1.5 · Default = 1.0",
        "Location Weight: Named junction = 2.0 · Main-road violation (no junction) = 1.5 · Other = 1.0",
        "Example: HGV + Main road parking + Named junction = 5.0 × 3.0 × 2.0 = 30.0 PCIS",
        "Example: Scooter + Wrong parking + Off-grid = 1.0 × 2.0 × 1.0 = 2.0 PCIS",
    ])

    add_content_slide(prs, "A2 — Road Profiles & Flow Assumptions", [
        "Main road: 10.5 m width · 3 lanes · 1,200 PCU/lane/hr · IRC:86 urban arterial",
        "Junction approach: 10.5 m · 3 lanes · 1,000 PCU/lane/hr",
        "Residential: 6.0 m · 2 lanes · 800 PCU/lane/hr",
        "Parking duration: 2 hours average per violation event",
        "Per-vehicle delay at obstruction: 8 seconds (queuing/merging)",
        "Vehicle occupancy: 1.5 persons per vehicle (BMTC urban survey reference)",
        "Fine amount: ₹500 per violation (MV Act standard challan)",
        "Capacity reduction: geometric width-based lower bound — real loss typically worse",
    ])

    add_content_slide(prs, "A3 — Replay Configuration", [
        "Holdout days: 14 (2024-03-26 to 2024-04-08)",
        "Patrol units: 5 · Shift: 4 hours · Dwell: 30 min/stop → 8 stops/unit → 40 total stops",
        "Intercept radius: 0.35 km (same ACTUAL and PRAHARI)",
        "ACTUAL strategy: top junctions by training-period violation count (fixed, not shift-aware)",
        "PRAHARI junction: causal PCIS ranking per shift window, junction names only",
        "PRAHARI hex: causal PCIS ranking per shift window, H3 resolution-9 cells (includes off-grid)",
        "Scoring metric: PCIS (not raw count) — weights severity, not just frequency",
        "Shift windows scored: 08–12, 12–16, 16–20, 20–24 IST",
    ])

    add_content_slide(prs, "A4 — Model Features (Final)", [
        "Included: time_window · dow · is_weekend · month · hour_sin · hour_cos · dow_sin · dow_cos",
        "  lag_1 · lag_6 · lag_42 · rolling_6_mean · rolling_42_mean · recurrence_factor",
        "Excluded (leakage): heavy_ratio · main_road_ratio · any post-aggregation composition ratio",
        "Train cutoff: 2024-02-15 · Test: after cutoff",
        "Classifier: LightGBM 300 trees · max_depth 6 · lr 0.05",
        "Regressor: same hyperparameters on positive samples only",
        "Evaluation: Precision@K · Hit Rate@K · MAE · RMSE · F1",
    ])

    add_content_slide(prs, "A5 — Repository Structure", [
        "app.py — Streamlit Command Center (entry point)",
        "config.py — paths, weights, defaults (all relative, no hardcoded local paths)",
        "src/data_pipeline.py — load, clean, PCIS, aggregate",
        "src/spatial_analysis.py — HDBSCAN, exposure, displacement",
        "src/prediction.py — hurdle model",
        "src/traffic_flow_impact.py — IRC/HCM delay model",
        "src/optimizer.py — VRP + greedy routing + scenario sim",
        "src/enforcement_replay.py — held-out validation",
        "src/patrol_briefing.py — field card generation",
        "scripts/bundle_data.py — one-time cloud deployment bundler",
        "data/*.parquet + flow_summary.json + default_patrol_plan.json",
    ])

    add_content_slide(prs, "A6 — Headline Numbers (Verified)", [
        f"Total parking violations: {total:,}",
        f"Off-grid (no junction code): {off_grid:,} ({off_pct:.1f}%)",
        f"Hidden HDBSCAN clusters: 1,436",
        f"Daily vehicle-hours delay (est.): ~{flow.get('daily_avg_delay_veh_hours', 5075):,.0f}",
        f"Peak hour (IST): {s.get('peak_hour', 10)}:00",
        f"Replay junction: +{rj.get('improvement_pct', 1.4):.1f}% PCIS",
        f"Replay hex: +{rh.get('improvement_pct', 5.2):.1f}% PCIS",
        f"Model Precision@10%: {mm.get('precision_at_10pct', 1.0)*100:.0f}%",
        f"Default plan PCIS intercepted: {plan.get('total_intercepted_pcis', 273228):,.0f}",
        f"Bundled data size: ~24 MB · Raw CSV: 105 MB (not committed)",
    ])

    # ── CLOSING ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_white(slide)
    t = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(8.4), Inches(1.0))
    t.text_frame.paragraphs[0].text = "PRAHARI"
    t.text_frame.paragraphs[0].font.size = Pt(40)
    t.text_frame.paragraphs[0].font.bold = True
    t.text_frame.paragraphs[0].font.color.rgb = ACCENT

    c = slide.shapes.add_textbox(Inches(0.8), Inches(3.2), Inches(8.4), Inches(2.0))
    cf = c.text_frame
    closing = [
        "BTP cannot enforce what it cannot see.",
        f"PRAHARI sees all {total:,} violations — including {off_grid:,} the junction system misses.",
        "Shift-ready patrol deployment. Held-out proof. Transparent assumptions.",
        "",
        "Live demo: [Streamlit URL — add after deployment]",
        "Repository: [GitHub URL — add before submission]",
        "",
        "Thank you.",
    ]
    for i, line in enumerate(closing):
        p = cf.paragraphs[0] if i == 0 else cf.add_paragraph()
        p.text = line
        p.font.size = Pt(16 if i < 4 else 14)
        p.font.color.rgb = DARK if i < 4 else GRAY
        if i == 0:
            p.font.bold = True

    prs.save(OUTPUT)
    print(f"Saved {OUTPUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build_ppt()
